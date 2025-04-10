# module2/data_parsing.py
import re
import unicodedata
from pathlib import Path

import pdfplumber
from pptx import Presentation
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, Settings, get_response_synthesizer
from llama_index.core.ingestion import IngestionPipeline
from llama_index.vector_stores.chroma import ChromaVectorStore
from chromadb import PersistentClient
from llama_index.core.schema import Document

from module2.config import Config
from module2.utils import setup_settings, ensure_persist_dir

# CUSTOM PPTX READER
from llama_index.core.readers.base import BaseReader

class PPTXTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        cleaned_text = preprocess_text(file_path, text)
                        text_runs.append(cleaned_text)
        full_text = "\n".join(text_runs)
        return [Document(text=full_text, metadata={"file_path": str(file_path)})]

# CUSTOM PDF READER
class PDFTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        text_runs = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    cleaned_text = preprocess_text(file_path, text)
                    text_runs.append(cleaned_text)
        full_text = "\n".join(text_runs)
        return [Document(text=full_text, metadata={"file_path": str(file_path)})]

# TEXT PREPROCESSING FUNCTIONS
def preprocess_text(file_path, text):
    ext = Path(file_path).suffix.lower()
    text_extensions = {".pdf", ".docx", ".pptx", ".txt"}
    if ext in text_extensions:
        text = clean_text(text)
    return text

def clean_text(text):
    text = unicodedata.normalize("NFKD", text)
    # Preserve email addresses
    email_pattern = r'([A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})'
    emails = re.findall(email_pattern, text)
    for i, email in enumerate(emails):
        text = text.replace(email, f'EMAIL_PLACEHOLDER_{i}')
    # Preserve URLs
    url_pattern = r'http[s]?://\S+'
    urls = re.findall(url_pattern, text)
    for i, url in enumerate(urls):
        text = text.replace(url, f'URL_PLACEHOLDER_{i}')
    text = re.sub(r'\.{5,}', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    for i, email in enumerate(emails):
        text = text.replace(f'EMAIL_PLACEHOLDER_{i}', email)
    for i, url in enumerate(urls):
        text = text.replace(f'URL_PLACEHOLDER_{i}', url)
    return text

def chunk_text(text, chunk_size=512, overlap=50):
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end]
        chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

def main():
    setup_settings(Config)

    # Use pathlib for robust path handling
    doc_path = Path("data/Course_Modules")
    if not doc_path.exists():
        print(f"❌ Path {doc_path} does not exist.")
        return

    try:
        reader = SimpleDirectoryReader(
            input_dir=str(doc_path),
            required_exts=['.pptx', '.ipynb', '.docx', '.csv', '.jpeg', '.pdf', '.png', '.py'],
            file_extractor={
                ".pptx": PPTXTextOnlyReader(),
                ".pdf": PDFTextOnlyReader()
            },
            recursive=True
        )
        docs = reader.load_data()
        if not docs:
            print("❌ No documents found! Check the path and file extensions.")
            return
        print(f"✅ Loaded {len(docs)} docs")
    except Exception as e:
        print(f"❌ Error loading documents: {e}")
        return

    processed_docs = []
    CHUNK_SIZE = 512
    CHUNK_OVERLAP = 50

    for doc in docs:
        doc_text = preprocess_text(doc.metadata.get("file_path", ""), doc.get_content())
        text_chunks = chunk_text(doc_text, CHUNK_SIZE, CHUNK_OVERLAP)
        folder_name = Path(doc.metadata.get("file_path", "")).parent.name
        for chunk in text_chunks:
            metadata = {
                "file_name": str(doc.metadata.get("file_name", "")),
                "file_path": str(doc.metadata.get("file_path", "")),
                "folder_name": folder_name,
                "num_tokens": len(chunk.split()),
                "num_chars": len(chunk),
            }
            processed_docs.append({
                "doc_id": doc.doc_id,
                "text": chunk,
                "metadata": metadata,
                "category": "<category>"
            })

    document_objects = [Document(text=doc["text"], metadata=doc["metadata"])
                        for doc in processed_docs]

    pipeline = IngestionPipeline(
        transformations=[Settings.embed_model]  # Applying embedding transformation
    )

    try:
        nodes = pipeline.run(documents=document_objects)
        if not nodes:
            print("❌ No nodes were created. Check document parsing.")
            return
        print(f"✅ {len(nodes)} document nodes created and stored in ChromaDB.")
        
        chroma_path = Path("./chroma_db")
        chroma_client = PersistentClient(path=str(chroma_path))
        collection = chroma_client.get_or_create_collection("document_chunks")
        vector_store = ChromaVectorStore(chroma_client, collection_name="document_chunks")
        for i, node in enumerate(nodes):
            collection.add(
                ids=[str(i)],
                documents=[node.text],
                metadatas=[node.metadata]
            )
    except Exception as e:
        print(f"❌ Error during ingestion pipeline: {e}")
        return

    try:
        index = VectorStoreIndex(nodes, vector_store=vector_store)
        print("✅ Vector store index created successfully.")
        persist_dir = ensure_persist_dir("persisted_index")
        index.storage_context.persist(persist_dir=str(persist_dir))
        print(f"✅ Index persisted to {persist_dir}")
    except Exception as e:
        print(f"❌ Error creating VectorStoreIndex: {e}")
        return

    from module2.utils import RAGQueryEngine
    retriever = index.as_retriever()
    synthesizer = get_response_synthesizer(response_mode="compact")
    query_engine = RAGQueryEngine(retriever=retriever, response_synthesizer=synthesizer)
    response = query_engine.query("Write a code to find a factorial for a number?")
    print(response)

if __name__ == '__main__':
    main()
