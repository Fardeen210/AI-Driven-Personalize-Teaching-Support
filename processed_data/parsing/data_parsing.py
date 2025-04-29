# module2/data_parsing.py

import pickle
from pathlib import Path

import pdfplumber
from pptx import Presentation

from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import (
    SimpleDirectoryReader, Settings,
)
from llama_index.core.ingestion import IngestionPipeline
from llama_index.core.schema import Document

# Import configuration and utility functions.
from processed_data.module2.config import Config
from processed_data.module2.utils import setup_settings

# ------------------------------
# CUSTOM PPTX READER (NO PREPROCESSING)
# ------------------------------
from llama_index.core.readers.base import BaseReader

class PPTXTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        # Use raw text without additional preprocessing.
                        text_runs.append(text)
        full_text = "\n".join(text_runs)
        return [Document(text=full_text, metadata={"file_path": str(file_path)})]

# ------------------------------
# CUSTOM PDF READER (NO PREPROCESSING)
# ------------------------------
class PDFTextOnlyReader(BaseReader):
    def load_data(self, file_path: str, extra_info=None) -> list[Document]:
        text_runs = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    # Use raw text without additional preprocessing.
                    text_runs.append(text)
        full_text = "\n".join(text_runs)
        return [Document(text=full_text, metadata={"file_path": str(file_path)})]

# ------------------------------
# TEXT CHUNKING FUNCTION
# ------------------------------
def chunk_text(text, chunk_size=512, overlap=50):
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

# ------------------------------
# MAIN FUNCTION: PARSING & INGESTION
# ------------------------------
def main():
    # Set up configuration (including embedding model settings)
    setup_settings(Config)
    
    # Define the folder where your course modules reside.
    doc_path = Path("data/Course_Modules")
    if not doc_path.exists():
        print(f"❌ Path {doc_path} does not exist.")
        return

    # Load documents using SimpleDirectoryReader with custom extractors.
    try:
        reader = SimpleDirectoryReader(
            input_dir=str(doc_path),
            required_exts=['.pptx', '.ipynb', '.docx', '.csv', '.jpeg', '.pdf', '.png', '.py'],
            file_extractor={
                ".pptx": PPTXTextOnlyReader(),
                ".pdf": PDFTextOnlyReader()
            },
            recursive=True
        )
        docs = reader.load_data()
        if not docs:
            print("❌ No documents found! Check the path and file extensions.")
            return
        print(f"✅ Loaded {len(docs)} docs")
    except Exception as e:
        print(f"❌ Error loading documents: {e}")
        return

    processed_docs = []
    CHUNK_SIZE = 512
    CHUNK_OVERLAP = 50

    # Process each document by chunking text and preparing metadata.
    for doc in docs:
        # Use the raw text from the document.
        doc_text = doc.get_content()
        text_chunks = chunk_text(doc_text, CHUNK_SIZE, CHUNK_OVERLAP)
        folder_name = Path(doc.metadata.get("file_path", "")).parent.name
        for chunk in text_chunks:
            metadata = {
                "file_name": str(doc.metadata.get("file_name", "")),
                "file_path": str(doc.metadata.get("file_path", "")),
                "folder_name": folder_name,
                "num_tokens": len(chunk.split()),
                "num_chars": len(chunk),
            }
            processed_docs.append({
                "doc_id": doc.doc_id,
                "text": chunk,
                "metadata": metadata,
                "category": "<category>"  # Adjust as needed.
            })

    # Convert processed documents into Document objects.
    document_objects = [
        Document(text=doc_info["text"], metadata=doc_info["metadata"])
        for doc_info in processed_docs
    ]

    # Run the ingestion pipeline using the embedding transformation.
    pipeline = IngestionPipeline(transformations=[Settings.embed_model])
    try:
        nodes = pipeline.run(documents=document_objects)
        if not nodes:
            print("❌ No nodes were created. Check document parsing.")
            return
        print(f"✅ {len(nodes)} document nodes created.")
    except Exception as e:
        print(f"❌ Error during ingestion pipeline: {e}")
        return

    # Save the nodes to an intermediate file (parsed_nodes.pkl)
    try:
        with open("parsed_nodes.pkl", "wb") as f:
            pickle.dump(nodes, f)
        print("✅ Parsed nodes saved to parsed_nodes.pkl")
    except Exception as e:
        print(f"❌ Error saving parsed nodes: {e}")

if __name__ == '__main__':
    main()